import streamlit as st
import pandas as pd
import numpy as np
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import re
from io import BytesIO
import subprocess, tempfile, os, json

# ─────────────────────────────────────────────
# PAGE CONFIG
# ─────────────────────────────────────────────
st.set_page_config(
    page_title="Whirlpool Production Plan Dashboard",
    page_icon="❄️",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
.metric-card {
    background: #f0f4ff;
    border-radius: 10px;
    padding: 16px 20px;
    border-left: 4px solid #1d4ed8;
    margin-bottom: 8px;
}
.metric-label { font-size: 12px; color: #6b7280; font-weight: 600; text-transform: uppercase; }
.metric-value { font-size: 28px; font-weight: 700; color: #1e293b; }
.metric-sub   { font-size: 12px; color: #6b7280; margin-top: 2px; }
.section-header {
    font-size: 18px; font-weight: 700; color: #1e293b;
    border-bottom: 2px solid #e2e8f0; padding-bottom: 6px; margin: 24px 0 12px 0;
}
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────
# PARSING
# ─────────────────────────────────────────────
def parse_release_date(raw):
    if pd.isna(raw):
        return "—"
    s = str(raw).strip()
    if hasattr(raw, 'strftime'):
        return raw.strftime("%d %b '%y")
    match = re.search(r'(\d{1,2})[a-z]{0,2}[\s\-]([A-Za-z]+)', s)
    if match:
        return f"{match.group(1)} {match.group(2)}"
    return s[:15]

@st.cache_data(show_spinner=False)
def parse_excel(file_bytes, filename):
    if isinstance(file_bytes, bytes):
        file_bytes = BytesIO(file_bytes)
    df_raw = pd.read_excel(file_bytes, sheet_name=0, header=None)

    date_row  = df_raw.iloc[17]
    total_row = df_raw.iloc[18]

    # --- Auto-detect version columns: row 17 cells containing "Released"
    version_cols = sorted(
        [c for c in range(len(date_row))
         if pd.notna(date_row.iloc[c]) and 'released' in str(date_row.iloc[c]).lower()],
        reverse=True   # highest index = oldest, lowest index = most recent
    )

    wk_map = {}

    # --- Latest version = lowest-index col with non-zero total (handles unreleased months)
    latest_col = min(version_cols)   # fallback: most recent = lowest col index
    for c in sorted(version_cols):
        if (pd.to_numeric(total_row.iloc[c], errors='coerce') or 0) > 0:
            latest_col = c
            break

    first_col = max(version_cols)    # oldest = highest col index

    # --- Version labels ordered oldest → newest
    version_labels = [parse_release_date(date_row.iloc[c]) for c in sorted(version_cols, reverse=True)]

    # --- Parse SKU rows
    raw = df_raw.iloc[19:].copy().reset_index(drop=True)
    raw.columns = range(raw.shape[1])

    data = pd.DataFrame()
    data['MAT']         = raw[0].astype(str).str.strip()
    data['Color']       = raw[1].astype(str).str.strip()
    data['Star']        = raw[2].astype(str).str.strip()
    data['Segment']     = raw[3].astype(str).str.strip()
    data['Description'] = raw[4].astype(str).str.strip()
    data['v1_final']    = pd.to_numeric(raw[latest_col], errors='coerce').fillna(0).astype(int)
    data['v5_first']    = pd.to_numeric(raw[first_col],  errors='coerce').fillna(0).astype(int)

    # Intermediate versions for plan evolution chart
    mid_cols = sorted([c for c in version_cols if c != latest_col and c != first_col])
    data['_version_cols'] = str(version_cols)   # store for evolution chart
    # Store all version totals as metadata on the df (via attrs)
    all_version_cols = sorted(version_cols, reverse=True)  # old → new



    # Store per-version columns for the evolution chart
    for i, c in enumerate(all_version_cols):
        data[f'ver_{i}'] = pd.to_numeric(raw[c], errors='coerce').fillna(0).astype(int)

    data = data[~data['MAT'].isin(['nan', 'NaN', '', 'None'])]

    data['Is_Active'] = data['v1_final'] > 0
    data['Is_TBC']    = data['MAT'].str.upper().str.startswith('TBC')
    data['drift_abs'] = data['v1_final'] - data['v5_first']
    data['drift_pct'] = np.where(
        data['v5_first'] > 0,
        (data['drift_abs'] / data['v5_first'] * 100).round(1),
        np.nan
    )

    def sku_status(row):
        f, l = row['v5_first'], row['v1_final']
        if f == 0 and l > 0:  return 'New'
        if f > 0 and l == 0:  return 'Dropped'
        if f == l:             return 'Unchanged'
        return 'Increased' if l > f else 'Decreased'

    data['Status'] = data.apply(sku_status, axis=1)
    return data, version_labels

# ─────────────────────────────────────────────
# CHART HELPERS
# ─────────────────────────────────────────────
BLUE   = '#1d4ed8'
GREEN  = '#16a34a'
RED    = '#dc2626'
AMBER  = '#d97706'
PURPLE = '#7c3aed'
GRAY   = '#6b7280'

COLOR_MAP = {'WINE':'#8B1A1A','BLUE':'#1d4ed8','GREY':'#9ca3af','BLACK':'#1f2937','PURPLE':'#7c3aed',
             'Wine':'#8B1A1A','Blue':'#1d4ed8','Grey':'#9ca3af','Black':'#1f2937','Purple':'#7c3aed'}
STATUS_COLOR = {'New':GREEN,'Increased':'#86efac','Unchanged':'#94a3b8','Decreased':'#fca5a5','Dropped':RED}

def fmt(n): return f"{int(n):,}"

def apply_gradient(styler, subset, col_data):
    """Light RdYlGn gradient with white text for readability."""
    import matplotlib.cm as cm
    import matplotlib.colors as mcolors
    cmap = cm.get_cmap('RdYlGn')
    def cell_style(val):
        if pd.isna(val): return 'background-color: #f9fafb; color: #6b7280'
        norm = min(max((val + 60) / 120, 0), 1)  # -60..+60 → 0..1
        r, g, b, _ = cmap(norm)
        # Lighten: blend 30% toward white for background
        r2 = r + (1-r)*0.30; g2 = g + (1-g)*0.30; b2 = b + (1-b)*0.30
        hex_bg = mcolors.to_hex((r2, g2, b2))
        # Text: darken the same hue by blending 40% toward black
        r3 = r * 0.45; g3 = g * 0.45; b3 = b * 0.45
        hex_text = mcolors.to_hex((r3, g3, b3))
        return f'background-color: {hex_bg}; color: {hex_text}; font-weight: 700'
    return styler.map(cell_style, subset=subset)

def plan_evolution_chart(df, version_labels):
    # Use dynamically stored ver_N columns (ver_0 = oldest, ver_N = newest)
    ver_cols = sorted([c for c in df.columns if str(c).startswith('ver_')],
                      key=lambda x: int(x.split('_')[1]))
    totals = [df[c].sum() for c in ver_cols]
    # Trim version_labels to match available columns
    labels = version_labels[-len(totals):]
    fig = go.Figure(go.Scatter(
        x=labels, y=totals, mode='lines+markers+text',
        text=[fmt(t) for t in totals], textposition='top center',
        line=dict(color=BLUE, width=3), marker=dict(size=10, color=BLUE),
        fill='tozeroy', fillcolor='rgba(29,78,216,0.07)',
    ))
    fig.update_layout(title=dict(text='Plan Evolution — Total Units Across All Versions', x=0, xanchor='left', font=dict(size=14)),
        xaxis_title='Version (chronological →)', yaxis_title='Total Units',
        yaxis_tickformat=',', height=360, margin=dict(t=60,b=50,l=70,r=30),
        plot_bgcolor='white', paper_bgcolor='white',
        xaxis=dict(showgrid=False, tickfont=dict(size=11)),
        yaxis=dict(showgrid=True, gridcolor='#f1f5f9'))
    return fig

def segment_drift_chart(df):
    seg = df.groupby('Segment').agg(first=('v5_first','sum'), final=('v1_final','sum')).reset_index()
    seg = seg[~seg['Segment'].isin(['Segment','nan','NaN','','None'])]
    seg = seg[(seg['first'] > 0) | (seg['final'] > 0)]
    seg['drift_pct'] = np.where(seg['first'] > 0, ((seg['final']-seg['first'])/seg['first']*100).round(1), np.nan)
    seg = seg.dropna(subset=['drift_pct']).sort_values('drift_pct')
    colors = [GREEN if x >= 0 else RED for x in seg['drift_pct']]
    fig = go.Figure(go.Bar(
        x=seg['drift_pct'], y=seg['Segment'], orientation='h', marker_color=colors,
        text=[f"{x:+.1f}%" for x in seg['drift_pct']], textposition='outside',
    ))
    fig.add_vline(x=0, line_color='#374151', line_width=1.5)
    fig.update_layout(title=dict(text='Segment Drift — 1st Tentative → Final (%)', x=0, xanchor='left', font=dict(size=14)),
        xaxis_title='% Change', height=max(340, len(seg)*42),
        margin=dict(t=60,b=50,l=150,r=100),
        plot_bgcolor='white', paper_bgcolor='white',
        xaxis=dict(showgrid=True, gridcolor='#f1f5f9'),
        yaxis=dict(showgrid=False, tickfont=dict(size=11)))
    return fig

def color_mix_donut(df):
    active = df[df['Is_Active']]
    ca = active.groupby('Color')['v1_final'].sum().reset_index()
    ca = ca[ca['v1_final'] > 0].sort_values('v1_final', ascending=False)
    total = ca['v1_final'].sum()
    # Normalize color keys to uppercase for lookup
    colors = [COLOR_MAP.get(c.upper(), COLOR_MAP.get(c, '#94a3b8')) for c in ca['Color']]
    fig = go.Figure(go.Pie(
        labels=ca['Color'], values=ca['v1_final'], hole=0.55,
        marker_colors=colors,
        textinfo='label+percent', textfont_size=12,
    ))
    fig.update_layout(title='Color Mix — Final Plan',
        annotations=[dict(text=f"<b>{fmt(total)}</b><br>units", x=0.5, y=0.5, font_size=14, showarrow=False)],
        height=380, margin=dict(t=55,b=20,l=20,r=120), paper_bgcolor='white',
        legend=dict(orientation='v', x=1.02, y=0.5))
    return fig

def star_rating_chart(df):
    active = df[df['Is_Active']]
    sa = active.groupby('Star')['v1_final'].sum().reset_index()
    sa = sa[sa['v1_final'] > 0].sort_values('Star')
    fig = go.Figure(go.Bar(
        x=sa['Star'], y=sa['v1_final'], marker_color=PURPLE,
        text=[fmt(v) for v in sa['v1_final']], textposition='outside',
    ))
    fig.update_layout(title=dict(text='Star Rating Mix — Final Plan', x=0, xanchor='left', font=dict(size=14)),
        xaxis_title='Star Rating', yaxis_title='Units', yaxis_tickformat=',',
        height=320, margin=dict(t=60,b=50,l=70,r=30),
        plot_bgcolor='white', paper_bgcolor='white',
        xaxis=dict(showgrid=False), yaxis=dict(showgrid=True, gridcolor='#f1f5f9'))
    return fig



def sku_status_chart(df):
    order  = ['New','Increased','Unchanged','Decreased','Dropped']
    counts = df['Status'].value_counts().reindex(order, fill_value=0)
    fig = go.Figure(go.Bar(
        x=counts.index, y=counts.values,
        marker_color=[STATUS_COLOR[s] for s in order],
        text=counts.values, textposition='outside',
    ))
    fig.update_layout(title=dict(text='SKU Status — 1st Tentative → Final', x=0, xanchor='left', font=dict(size=14)),
        yaxis_title='# SKUs', height=320, margin=dict(t=60,b=50,l=60,r=30),
        plot_bgcolor='white', paper_bgcolor='white',
        xaxis=dict(showgrid=False), yaxis=dict(showgrid=True, gridcolor='#f1f5f9'))
    return fig

def segment_breakdown_chart(df):
    seg = df.groupby('Segment').agg(first=('v5_first','sum'), final=('v1_final','sum')).reset_index()
    seg = seg[~seg['Segment'].isin(['Segment','nan','NaN','','None'])]
    seg = seg[(seg['first'] > 0) | (seg['final'] > 0)].sort_values('final', ascending=False)
    fig = go.Figure()
    fig.add_trace(go.Bar(name='1st Tentative', x=seg['Segment'], y=seg['first'], marker_color='#93c5fd'))
    fig.add_trace(go.Bar(name='Final Plan',    x=seg['Segment'], y=seg['final'], marker_color=BLUE))
    fig.update_layout(barmode='group',
        title=dict(text='Segment Volume: 1st Tentative vs Final', x=0, xanchor='left', font=dict(size=14)),
        yaxis_title='Units', yaxis_tickformat=',',
        height=380, margin=dict(t=100,b=90,l=70,r=20),
        plot_bgcolor='white', paper_bgcolor='white',
        xaxis=dict(tickangle=-40, showgrid=False, tickfont=dict(size=11)),
        yaxis=dict(showgrid=True, gridcolor='#f1f5f9'),
        legend=dict(orientation='h', x=0, y=1.08))
    return fig

def multi_month_overview_chart(all_data):
    rows = []
    for month, (df, _) in all_data.items():
        first = df['v5_first'].sum()
        final = df['v1_final'].sum()
        drift = round((final-first)/first*100, 1) if first > 0 else 0
        rows.append({'Month': month, 'First': first, 'Final': final, 'Drift': drift})
    mdf = pd.DataFrame(rows)
    fig = make_subplots(rows=1, cols=2,
        subplot_titles=('Total Volume by Month', 'Overall Drift % vs 1st Tentative'))
    fig.add_trace(go.Bar(x=mdf['Month'], y=mdf['First'], name='1st Tentative', marker_color='#93c5fd'), row=1, col=1)
    fig.add_trace(go.Bar(x=mdf['Month'], y=mdf['Final'], name='Final', marker_color=BLUE), row=1, col=1)
    fig.add_trace(go.Bar(x=mdf['Month'], y=mdf['Drift'],
        marker_color=[GREEN if d >= 0 else RED for d in mdf['Drift']],
        text=[f"{d:+.1f}%" for d in mdf['Drift']], textposition='outside',
        showlegend=False), row=1, col=2)
    fig.add_hline(y=0, line_color='#374151', line_width=1, row=1, col=2)
    fig.update_layout(height=380, barmode='group', paper_bgcolor='white',
        showlegend=True, margin=dict(t=80,b=60,l=70,r=30))
    fig.update_yaxes(tickformat=',', row=1, col=1, showgrid=True, gridcolor='#f1f5f9')
    fig.update_yaxes(ticksuffix='%', row=1, col=2, showgrid=True, gridcolor='#f1f5f9')
    return fig

def multi_month_segment_chart(all_data):
    """Grouped bar: segments on x-axis, one bar per month (final plan qty)"""
    months = list(all_data.keys())
    all_segs = sorted(set(
        seg for df, _ in all_data.values()
        for seg in df['Segment'].unique()
        if seg not in ('Segment', 'nan', 'NaN', '', 'None')
        and (df[df['Segment']==seg]['v1_final'].sum() > 0 or df[df['Segment']==seg]['v5_first'].sum() > 0)
    ))
    palette = ['#1d4ed8','#16a34a','#d97706','#dc2626','#7c3aed','#0891b2','#be185d']
    fig = go.Figure()
    for i, month in enumerate(months):
        df = all_data[month][0]
        seg_finals = [df[df['Segment']==s]['v1_final'].sum() for s in all_segs]
        fig.add_trace(go.Bar(
            name=month, x=all_segs, y=seg_finals,
            marker_color=palette[i % len(palette)],
            text=[fmt(v) if v > 0 else '' for v in seg_finals],
            textposition='outside', textfont=dict(size=8),
        ))
    fig.update_layout(
        title=dict(text='Final Plan Volume by Segment — All Months', x=0, xanchor='left', font=dict(size=14)),
        barmode='group', height=460,
        xaxis=dict(tickangle=-40, showgrid=False, tickfont=dict(size=11)),
        yaxis=dict(tickformat=',', showgrid=True, gridcolor='#f1f5f9'),
        plot_bgcolor='white', paper_bgcolor='white',
        margin=dict(t=100,b=100,l=60,r=20),
        legend=dict(orientation='h', x=0, y=1.08, font=dict(size=11))
    )
    return fig

def multi_month_drift_heatmap(all_data):
    """Heatmap: rows=segments, cols=months, values=drift%"""
    months = list(all_data.keys())
    all_segs = sorted(set(
        seg for df, _ in all_data.values()
        for seg in df['Segment'].unique()
        if seg not in ('Segment', 'nan', 'NaN', '', 'None')
        and (df[df['Segment']==seg]['v1_final'].sum() > 0 or df[df['Segment']==seg]['v5_first'].sum() > 0)
    ))
    z, text = [], []
    for seg in all_segs:
        z_row, t_row = [], []
        for month in months:
            df = all_data[month][0]
            first = df[df['Segment']==seg]['v5_first'].sum()
            final = df[df['Segment']==seg]['v1_final'].sum()
            if first > 0:
                dp = round((final-first)/first*100, 1)
                z_row.append(dp)
                t_row.append(f"{dp:+.1f}%")
            else:
                z_row.append(None)
                t_row.append('—')
        z.append(z_row)
        text.append(t_row)
    fig = go.Figure(go.Heatmap(
        z=z, x=months, y=all_segs,
        text=text, texttemplate='%{text}',
        colorscale='RdYlGn', zmid=0,
        zmin=-60, zmax=60,
        colorbar=dict(title='Drift %', ticksuffix='%'),
    ))
    fig.update_layout(
        title=dict(text='Segment Drift % Heatmap — All Months (vs 1st Tentative)', x=0, xanchor='left', font=dict(size=14)),
        height=max(350, len(all_segs)*42 + 120),
        margin=dict(t=60,b=80,l=140,r=100),
        paper_bgcolor='white',
        xaxis=dict(side='bottom', tickfont=dict(size=11)),
        yaxis=dict(tickfont=dict(size=11)),
    )
    return fig

# ─────────────────────────────────────────────
# PPTX BUILDER (pure python-pptx, no subprocess)
# ─────────────────────────────────────────────
def build_pptx(df, version_labels, filename):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    import io, numpy as _np

    def rgb(h):
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def txb(slide, text, x, y, w, h, size=12, bold=False, color="1E293B",
            align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p  = tf.paragraphs[0]
        p.alignment = align
        r  = p.add_run()
        r.text = str(text)
        r.font.size    = Pt(size)
        r.font.bold    = bold
        r.font.italic  = italic
        r.font.color.rgb = rgb(color)
        r.font.name    = font
        return tb

    def rect(slide, x, y, w, h, fill_hex, line_hex=None, line_w=0):
        from pptx.util import Pt as _Pt
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_hex)
        if line_hex:
            shape.line.color.rgb = rgb(line_hex)
            shape.line.width = Emu(int(line_w * 12700))
        else:
            shape.line.fill.background()
        return shape

    # ── derive summary data
    ver_cols  = sorted([c for c in df.columns if str(c).startswith('ver_')],
                       key=lambda x: int(x.split('_')[1]))
    ver_totals = [int(df[c].sum()) for c in ver_cols]
    labels     = version_labels[-len(ver_totals):]

    total_final = int(df['v1_final'].sum())
    total_first = int(df['v5_first'].sum())
    drift_pct   = round((total_final-total_first)/total_first*100,1) if total_first else 0
    drift_abs   = total_final - total_first
    n_active    = int(df['Is_Active'].sum())
    sc          = df['Status'].value_counts()
    n_new       = int(sc.get('New',0))
    n_dropped   = int(sc.get('Dropped',0))
    n_increased = int(sc.get('Increased',0))
    n_decreased = int(sc.get('Decreased',0))
    n_unchanged = int(sc.get('Unchanged',0))

    seg = df.groupby('Segment').agg(first=('v5_first','sum'),final=('v1_final','sum')).reset_index()
    seg = seg[(seg['first']>0)|(seg['final']>0)].copy()
    seg['dp'] = _np.where(seg['first']>0,
        ((seg['final']-seg['first'])/seg['first']*100).round(1), _np.nan)
    seg = seg.dropna(subset=['dp'])

    color_agg = df[df['Is_Active']].groupby('Color')['v1_final'].sum().sort_values(ascending=False)
    month = filename.replace('.xlsx','').replace('.xls','').replace('_',' ')

    # ── Presentation setup
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # completely blank

    W, H = 13.33, 7.5
    NAVY="0D1F4E"; NAVYM="1A3370"; GOLD="D4A017"
    BLUE="1D4ED8"; WHITE="FFFFFF"; OFFWHITE="F4F6FA"
    MUTED="64748B"; GREEN="16A34A"; RED="DC2626"; AMBER="D97706"
    TEXT="1E293B"; LIGHT="EEF2FF"

    def sign(n): return "+" if n>=0 else ""
    def fmt(n):
        try: return f"{int(n):,}"
        except: return str(n)

    # ══════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, 0.3, H, GOLD)

    # Logo box
    rect(s, 0.55, 0.45, 3.0, 0.85, WHITE)
    txb(s,"Whirlpool",0.55,0.45,3.0,0.85,size=22,bold=True,color=NAVY,align=PP_ALIGN.CENTER,font="Georgia")

    txb(s,"Production Plan Review",0.55,1.55,9,1.1,size=44,bold=True,color=WHITE,font="Georgia")
    txb(s,month.upper(),0.55,2.75,5,0.5,size=22,color=GOLD,font="Calibri")
    txb(s,f"Plan versions: {'  →  '.join(labels)}",0.55,3.55,11,0.4,size=12,color="8BAEE0")

    kpis = [
        ("FINAL VOLUME",       fmt(total_final),           "units",    GOLD),
        ("DRIFT VS 1ST PLAN",  f"{sign(drift_pct)}{drift_pct}%", f"{sign(drift_abs)}{fmt(drift_abs)} units", GREEN if drift_pct>=0 else RED),
        ("ACTIVE SKUs",        str(n_active),              "in final plan", GOLD),
        ("NEW SKUs",           str(n_new),                 "added",    GREEN),
        ("DROPPED SKUs",       str(n_dropped),             "removed",  RED),
    ]
    kw, kg, kx0 = 2.25, 0.2, 0.55
    for i,(lbl,val,sub,ac) in enumerate(kpis):
        x = kx0 + i*(kw+kg)
        rect(s, x, 4.55, kw, 2.4, NAVYM)
        rect(s, x, 4.55, 0.07, 2.4, ac)
        txb(s, lbl,      x+0.15, 4.68, kw-0.2, 0.32, size=8,  bold=True,  color="8BAEE0")
        txb(s, val,      x+0.15, 5.05, kw-0.2, 0.95, size=24, bold=True,  color=ac, font="Georgia")
        txb(s, sub,      x+0.15, 6.05, kw-0.2, 0.35, size=11, color="8BAEE0")

    txb(s,"Whirlpool India  |  Purchase Department  |  MBA SIP",0.55,7.18,10,0.25,size=10,color="4A6090")

    # ══════════════════════════════════════════
    # SLIDE 2 — Plan Evolution
    # ══════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, OFFWHITE)
    rect(s, 0, 0, W, 1.05, NAVY)
    rect(s, 0, 0, 0.3, H, GOLD)
    txb(s,"Plan Evolution",0.5,0.15,8,0.75,size=28,bold=True,color=WHITE,font="Georgia")
    txb(s,f"{month}  ·  {len(labels)} versions",0.5,0.18,12.5,0.7,size=12,color=GOLD,align=PP_ALIGN.RIGHT)

    # Line chart
    cd = ChartData()
    cd.categories = labels
    cd.add_series("Total Units", ver_totals)
    ch = s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.45),Inches(1.25),Inches(8.4),Inches(5.7), cd).chart

    # Version table
    tx, ty = 9.1, 1.25
    txb(s,"VERSION SUMMARY",tx,ty,4.0,0.35,size=8,bold=True,color=MUTED)
    non_zero = [(l,v) for l,v in zip(labels,ver_totals) if v>0]
    ty2 = ty+0.38
    for i,(l,v) in enumerate(non_zero):
        bg = WHITE if i%2==0 else LIGHT
        rect(s, tx, ty2, 4.0, 0.5, bg)
        txb(s, l,      tx+0.1, ty2+0.05, 1.8, 0.4, size=11, color=TEXT)
        txb(s, fmt(v), tx+2.0, ty2+0.05, 1.9, 0.4, size=11, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
        ty2 += 0.52

    # Drift callout
    dc = GREEN if drift_pct>=0 else RED
    db = "DCFCE7" if drift_pct>=0 else "FEE2E2"
    rect(s, tx, ty2+0.25, 4.0, 1.3, db)
    txb(s,"TOTAL DRIFT",  tx+0.15, ty2+0.33, 3.7, 0.28, size=8, bold=True, color=MUTED)
    txb(s,f"{sign(drift_pct)}{drift_pct}%", tx+0.15, ty2+0.62, 2.2, 0.75, size=34, bold=True, color=dc, font="Georgia")
    txb(s,f"{sign(drift_abs)}{fmt(drift_abs)} units", tx+2.4, ty2+0.9, 1.5, 0.4, size=12, color=dc, align=PP_ALIGN.RIGHT)

    # ══════════════════════════════════════════
    # SLIDE 3 — Segment Drift
    # ══════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, OFFWHITE)
    rect(s, 0, 0, W, 1.05, NAVY)
    rect(s, 0, 0, 0.3, H, GOLD)
    txb(s,"Segment Analysis",0.5,0.15,8,0.75,size=28,bold=True,color=WHITE,font="Georgia")
    txb(s,"1st Tentative → Final Plan",0.5,0.18,12.5,0.7,size=12,color=GOLD,align=PP_ALIGN.RIGHT)

    seg_s = seg.sort_values('final', ascending=False)
    cd2 = ChartData()
    cd2.categories = list(seg_s['Segment'])
    cd2.add_series("1st Tentative", [int(v) for v in seg_s['first']])
    cd2.add_series("Final Plan",    [int(v) for v in seg_s['final']])
    ch2 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.45),Inches(1.25),Inches(8.3),Inches(5.7), cd2).chart
    ch2.series[0].format.fill.solid(); ch2.series[0].format.fill.fore_color.rgb = rgb("93C5FD")
    ch2.series[1].format.fill.solid(); ch2.series[1].format.fill.fore_color.rgb = rgb(BLUE)
    ch2.has_legend = True

    # Drift table
    tx, ty0 = 8.9, 1.25
    txb(s,"DRIFT SUMMARY",tx,ty0,4.2,0.35,size=8,bold=True,color=MUTED)
    rect(s, tx, ty0+0.38, 4.2, 0.4, NAVY)
    txb(s,"Segment", tx+0.1,ty0+0.38,2.1,0.4,size=10,bold=True,color=WHITE)
    txb(s,"Drift %", tx+2.3,ty0+0.38,1.8,0.4,size=10,bold=True,color=WHITE,align=PP_ALIGN.RIGHT)
    ty3 = ty0+0.8
    seg_sorted = seg.sort_values('dp')
    rh = min(0.36, 5.8/max(len(seg_sorted),1))
    for i,(_,row) in enumerate(seg_sorted.iterrows()):
        if ty3+rh > 7.35: break
        bg = WHITE if i%2==0 else LIGHT
        dc2 = GREEN if row['dp']>=0 else RED
        rect(s, tx, ty3, 4.2, rh, bg)
        txb(s, row['Segment'], tx+0.1,ty3,2.1,rh,size=10,color=TEXT)
        txb(s, f"{sign(row['dp'])}{row['dp']}%", tx+2.3,ty3,1.8,rh,size=10,bold=True,color=dc2,align=PP_ALIGN.RIGHT)
        ty3 += rh+0.02

    # ══════════════════════════════════════════
    # SLIDE 4 — Mix Analysis
    # ══════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, OFFWHITE)
    rect(s, 0, 0, W, 1.05, NAVY)
    rect(s, 0, 0, 0.3, H, GOLD)
    txb(s,"Product Mix",0.5,0.15,8,0.75,size=28,bold=True,color=WHITE,font="Georgia")
    txb(s,"Color Mix & SKU Status  ·  Final Plan",0.5,0.18,12.5,0.7,size=12,color=GOLD,align=PP_ALIGN.RIGHT)

    # Color donut
    cd3 = ChartData()
    cd3.categories = list(color_agg.index)
    cd3.add_series("Units", [int(v) for v in color_agg.values])
    ch3 = s.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT,
        Inches(0.45),Inches(1.2),Inches(5.5),Inches(5.6), cd3).chart
    color_hex = {"WINE":"8B1A1A","BLUE":BLUE,"GREY":"9CA3AF","BLACK":"1F2937","PURPLE":"7C3AED"}
    for i,col_name in enumerate(color_agg.index):
        if i < len(ch3.series[0].points):
            pt = ch3.series[0].points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = rgb(color_hex.get(col_name,"888888"))
    ch3.has_legend = True
    txb(s,"COLOR MIX",0.45,1.12,5.5,0.28,size=8,bold=True,color=MUTED,align=PP_ALIGN.CENTER)

    # SKU status cards
    statuses = [
        ("New",       n_new,       "DCFCE7", GREEN),
        ("Increased", n_increased, "DCFCE7", "15803D"),
        ("Unchanged", n_unchanged, "F1F5F9", MUTED),
        ("Decreased", n_decreased, "FEF3C7", AMBER),
        ("Dropped",   n_dropped,   "FEE2E2", RED),
    ]
    cw, cg, cx0 = 1.35, 0.17, 6.2
    txb(s,"SKU STATUS BREAKDOWN",cx0,1.12,(cw+cg)*5,0.28,size=8,bold=True,color=MUTED,align=PP_ALIGN.CENTER)
    for i,(lbl,val,bg,ac) in enumerate(statuses):
        x = cx0 + i*(cw+cg)
        rect(s, x, 1.42, cw, 1.75, bg)
        txb(s, str(val), x, 1.52, cw, 0.85, size=30, bold=True, color=ac, font="Georgia", align=PP_ALIGN.CENTER)
        txb(s, lbl,      x, 2.4,  cw, 0.35, size=10, bold=True, color=ac, align=PP_ALIGN.CENTER)

    # Color table
    tx4, ty4 = 6.2, 3.45
    txb(s,"COLOR BREAKDOWN",tx4,ty4,6.9,0.28,size=8,bold=True,color=MUTED)
    rect(s, tx4, ty4+0.3, 6.9, 0.38, NAVY)
    txb(s,"Color",  tx4+0.1,ty4+0.3,2.5,0.38,size=10,bold=True,color=WHITE)
    txb(s,"Units",  tx4+2.7,ty4+0.3,2.1,0.38,size=10,bold=True,color=WHITE,align=PP_ALIGN.RIGHT)
    txb(s,"Share",  tx4+4.9,ty4+0.3,1.9,0.38,size=10,bold=True,color=WHITE,align=PP_ALIGN.RIGHT)
    ty5 = ty4+0.7
    for i,(col_name,qty) in enumerate(color_agg.items()):
        if ty5+0.38 > 7.3: break
        bg = WHITE if i%2==0 else LIGHT
        pct = round(qty/total_final*100,1) if total_final else 0
        rect(s, tx4, ty5, 6.9, 0.38, bg)
        txb(s, col_name, tx4+0.1,ty5,2.5,0.38,size=11,color=TEXT)
        txb(s, fmt(qty),  tx4+2.7,ty5,2.1,0.38,size=11,bold=True,color=BLUE,align=PP_ALIGN.RIGHT)
        txb(s, f"{pct}%", tx4+4.9,ty5,1.9,0.38,size=11,color=MUTED,align=PP_ALIGN.RIGHT)
        ty5 += 0.4

    # ══════════════════════════════════════════
    # SLIDE 5 — Key Takeaways
    # ══════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 0, 0.3, H, GOLD)
    txb(s,"Key Takeaways",0.6,0.4,10,0.9,size=40,bold=True,color=WHITE,font="Georgia")
    txb(s,month.upper(),0.6,1.38,5,0.42,size=16,color=GOLD)

    top_gain = seg.sort_values('dp',ascending=False).iloc[0] if len(seg) else None
    top_lose = seg.sort_values('dp').iloc[0] if len(seg) else None
    top_col      = color_agg.index[0] if len(color_agg) else "N/A"
    top_col_qty  = int(color_agg.iloc[0]) if len(color_agg) else 0
    top_col_pct  = round(top_col_qty/total_final*100) if total_final else 0


    bullets = [
        f"Final plan: {fmt(total_final)} units  ({sign(drift_pct)}{drift_pct}% vs 1st tentative of {fmt(total_first)}).",
        f"Highest segment uptick: {top_gain['Segment']} ({sign(top_gain['dp'])}{top_gain['dp']}%) — strongest demand signal." if top_gain is not None else "",
        f"Biggest cut: {top_lose['Segment']} ({top_lose['dp']}%) — may warrant supply chain review." if top_lose is not None else "",
        f"Dominant color: {top_col} at {fmt(top_col_qty)} units ({top_col_pct}% of final plan).",
        f"{n_new} new SKUs introduced; {n_dropped} dropped — net mix churn of {n_new+n_dropped} SKUs.",
    ]
    for i, text in enumerate(bullets):
        if not text: continue
        y = 2.05 + i*0.77
        rect(s, 0.6, y+0.08, 0.3, 0.3, GOLD)
        txb(s, str(i+1), 0.6, y+0.05, 0.3, 0.35, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txb(s, text, 1.05, y, 11.9, 0.62, size=14, color=WHITE)

    rect(s, 0.6, 7.08, 12.4, 0.02, "2A4080")
    txb(s,"Whirlpool India  ·  Purchase Department  ·  MBA SIP  ·  Diwakar",0.6,7.15,12,0.28,size=10,color="4A6090")

    # ── Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ─────────────────────────────────────────────
# SIDEBAR — built once, fully
# ─────────────────────────────────────────────
with st.sidebar:
    st.image("whirlpool_logo.png", width=170)
    st.markdown("### Production Plan Dashboard")
    st.divider()

    uploaded_files = st.file_uploader(
        "Upload Plan File(s) (.xlsx)",
        type=['xlsx','xls'],
        accept_multiple_files=True,
    )

    # Parse files here so filters can be populated in the same sidebar block
    all_data = {}
    if uploaded_files:
        for f in uploaded_files:
            try:
                raw_bytes = f.read()
                st.session_state['_xlsx_bytes_' + f.name] = raw_bytes
                df_parsed, ver_labels = parse_excel(raw_bytes, f.name)
                all_data[f.name] = (df_parsed, ver_labels)
            except Exception as e:
                st.error(f"❌ {f.name}: {e}")

    st.divider()
    st.markdown("**Filters**")

    if all_data:
        all_dfs      = pd.concat([v[0] for v in all_data.values()], ignore_index=True)
        seg_options  = sorted(all_dfs['Segment'].dropna().unique().tolist())
        col_options  = sorted(all_dfs['Color'].dropna().unique().tolist())
        star_options = sorted(all_dfs['Star'].dropna().unique().tolist())
    else:
        seg_options = col_options = star_options = []

    filter_segments = st.multiselect("Segment",     options=seg_options)
    filter_colors   = st.multiselect("Color",       options=col_options)
    filter_stars    = st.multiselect("Star Rating", options=star_options)
    show_tbc        = st.checkbox("Include TBC SKUs", value=False)

    st.divider()
    if all_data:
        st.markdown("**Export**")
        sel_export = list(all_data.keys())[0]
        if st.button("📊 Generate PPTX Report", use_container_width=True):
            with st.spinner("Building presentation…"):
                try:
                    df_ex, vl_ex = all_data[sel_export]
                    pptx_bytes = build_pptx(df_ex, vl_ex, sel_export)
                    st.session_state['_pptx_bytes'] = pptx_bytes
                    st.session_state['_pptx_name']  = sel_export.replace('.xlsx','').replace('.xls','') + '_report.pptx'
                except Exception as e:
                    st.error(f"Export failed: {e}")
        if '_pptx_bytes' in st.session_state:
            st.download_button(
                "⬇️ Save PPTX",
                st.session_state['_pptx_bytes'],
                file_name=st.session_state.get('_pptx_name','report.pptx'),
                mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                use_container_width=True
            )
    st.divider()
    st.caption("Built with Streamlit + Plotly · Whirlpool India")

# ─────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────
st.markdown("## ❄️ Whirlpool Production Planning Dashboard")

if not uploaded_files or not all_data:
    st.info("👈 Upload one or more monthly plan `.xlsx` files using the sidebar to get started.")
    st.markdown("""
**Expected file format:**
- Standard Whirlpool production plan Excel (raw file, not the cleaned output)
- Rows 1–18: segment-level summary block · Row 17: release dates · Row 19: headers
- Rows 20+: SKU data — MAT | Color | Star | Segment | Description | v1_final → v5_first | Wk1–Wk4
    """)
    st.stop()

# Month selector
if len(all_data) == 1:
    selected_month = list(all_data.keys())[0]
else:
    selected_month = st.selectbox("📅 Select month to view:", list(all_data.keys()))

df, version_labels = all_data[selected_month]

# Apply filters
df_f = df.copy()
if filter_segments: df_f = df_f[df_f['Segment'].isin(filter_segments)]
if filter_colors:   df_f = df_f[df_f['Color'].isin(filter_colors)]
if filter_stars:    df_f = df_f[df_f['Star'].isin(filter_stars)]
if not show_tbc:    df_f = df_f[~df_f['Is_TBC']]

# ─────────────────────────────────────────────
# KPI STRIP
# ─────────────────────────────────────────────
active_df   = df_f[df_f['Is_Active']]
total_final = df_f['v1_final'].sum()
total_first = df_f['v5_first'].sum()
drift_total = total_final - total_first
drift_pct   = (drift_total / total_first * 100) if total_first else 0
n_active    = active_df['MAT'].nunique()
n_new       = (df_f['Status'] == 'New').sum()
n_dropped   = (df_f['Status'] == 'Dropped').sum()

c1, c2, c3, c4, c5 = st.columns(5)
with c1:
    st.markdown(f"""<div class='metric-card'>
        <div class='metric-label'>Final Plan Volume</div>
        <div class='metric-value'>{fmt(total_final)}</div>
        <div class='metric-sub'>units · latest version</div></div>""", unsafe_allow_html=True)
with c2:
    col = GREEN if drift_pct >= 0 else RED
    sign = '+' if drift_pct >= 0 else ''
    st.markdown(f"""<div class='metric-card' style='border-left-color:{col}'>
        <div class='metric-label'>Drift vs 1st Tentative</div>
        <div class='metric-value' style='color:{col}'>{sign}{drift_pct:.1f}%</div>
        <div class='metric-sub'>{sign}{fmt(drift_total)} units</div></div>""", unsafe_allow_html=True)
with c3:
    st.markdown(f"""<div class='metric-card'>
        <div class='metric-label'>Active SKUs</div>
        <div class='metric-value'>{n_active}</div>
        <div class='metric-sub'>in final plan</div></div>""", unsafe_allow_html=True)
with c4:
    st.markdown(f"""<div class='metric-card' style='border-left-color:{GREEN}'>
        <div class='metric-label'>New SKUs</div>
        <div class='metric-value'>{n_new}</div>
        <div class='metric-sub'>added vs 1st tentative</div></div>""", unsafe_allow_html=True)
with c5:
    st.markdown(f"""<div class='metric-card' style='border-left-color:{RED}'>
        <div class='metric-label'>Dropped SKUs</div>
        <div class='metric-value'>{n_dropped}</div>
        <div class='metric-sub'>removed vs 1st tentative</div></div>""", unsafe_allow_html=True)

st.caption(f"Version dates: {' → '.join(version_labels)}")

# ─────────────────────────────────────────────
# TABS
# ─────────────────────────────────────────────
tab1, tab2, tab3, tab4, tab5 = st.tabs([
    "📈 Plan Evolution", "📦 Segments", "🎨 Mix Analysis", "🔍 SKU Drill-Down", "📊 Multi-Month"
])

with tab1:
    st.markdown("<div class='section-header'>Plan Evolution</div>", unsafe_allow_html=True)
    st.plotly_chart(plan_evolution_chart(df_f, version_labels), use_container_width=True)
    ver_cols = sorted([c for c in df_f.columns if str(c).startswith('ver_')],
                      key=lambda x: int(x.split('_')[1]))
    totals   = [df_f[c].sum() for c in ver_cols]
    labels   = version_labels[-len(totals):]
    chgs     = ['—'] + [f"{((totals[i]-totals[i-1])/totals[i-1]*100):+.1f}%" if totals[i-1] else '—'
                        for i in range(1, len(totals))]
    st.dataframe(pd.DataFrame({'Version': labels, 'Total Units': [fmt(t) for t in totals],
        'Change vs Prev': chgs}), use_container_width=True, hide_index=True)

with tab2:
    st.markdown("<div class='section-header'>Segment Analysis</div>", unsafe_allow_html=True)
    ca, cb = st.columns(2)
    with ca: st.plotly_chart(segment_breakdown_chart(df_f), use_container_width=True)
    with cb: st.plotly_chart(segment_drift_chart(df_f), use_container_width=True)
    # Build segment table with ALL version columns
    ver_cols_sorted = sorted([c for c in df_f.columns if str(c).startswith('ver_')],
                              key=lambda x: int(x.split('_')[1]))
    ver_labels_trimmed = version_labels[-len(ver_cols_sorted):]

    # Display labels: first = "date / 1st Plan", last = "date / Final Plan", middle unchanged
    def make_display_label(lbl, i, total):
        if i == 0:          return f"{lbl} / 1st Plan"
        if i == total - 1:  return f"{lbl} / Final Plan"
        return lbl
    display_labels = [make_display_label(lbl, i, len(ver_labels_trimmed))
                      for i, lbl in enumerate(ver_labels_trimmed)]

    seg_agg = df_f.groupby('Segment').agg(Active_SKUs=('Is_Active','sum')).reset_index()
    for i, vc in enumerate(ver_cols_sorted):
        lbl = display_labels[i]
        seg_agg[lbl] = df_f.groupby('Segment')[vc].sum().values

    # Remove junk rows
    seg_agg = seg_agg[~seg_agg['Segment'].isin(['Segment','nan','NaN','','None'])]
    # Keep only rows with at least some data
    qty_cols = display_labels
    seg_agg = seg_agg[seg_agg[qty_cols].sum(axis=1) > 0]

    seg_agg['Drift']   = seg_agg[display_labels[-1]] - seg_agg[display_labels[0]]
    seg_agg['Drift %'] = np.where(seg_agg[display_labels[0]] > 0,
        (seg_agg['Drift'] / seg_agg[display_labels[0]] * 100).round(1), np.nan)

    fmt_cols = {lbl: '{:,.0f}' for lbl in display_labels}
    fmt_cols['Drift'] = '{:+,.0f}'
    fmt_cols['Drift %'] = '{:+.1f}%'

    # Use text-only gradient — no background_gradient on NaN cols to avoid black cells
    def drift_color(val):
        if pd.isna(val): return 'color: #6b7280'
        if val > 0:  return f'color: #15803d; font-weight:600'
        if val < 0:  return f'color: #dc2626; font-weight:600'
        return 'color: #6b7280'

    styled_seg2 = (
        seg_agg.sort_values(display_labels[-1], ascending=False)
        .style.format(fmt_cols, na_rep='—')
    )
    apply_gradient(styled_seg2, subset=['Drift %'], col_data=seg_agg['Drift %'])
    st.dataframe(styled_seg2, use_container_width=True, hide_index=True)

with tab3:
    st.markdown("<div class='section-header'>Mix Analysis — Final Plan</div>", unsafe_allow_html=True)
    ca, cb = st.columns(2)
    with ca: st.plotly_chart(color_mix_donut(df_f), use_container_width=True)
    with cb: st.plotly_chart(star_rating_chart(df_f), use_container_width=True)
    st.plotly_chart(sku_status_chart(df_f), use_container_width=True)

with tab4:
    st.markdown("<div class='section-header'>SKU-Level Data</div>", unsafe_allow_html=True)
    view_mode = st.radio("Show:", ["Active SKUs only","All SKUs","Dropped SKUs","New SKUs"], horizontal=True)
    drill = {'Active SKUs only': df_f[df_f['Is_Active']], 'All SKUs': df_f,
             'Dropped SKUs': df_f[df_f['Status']=='Dropped'],
             'New SKUs': df_f[df_f['Status']=='New']}[view_mode]
    show_cols = [c for c in ['MAT','Description','Segment','Color','Star',
        'v5_first','v4','v3','v2','v1_final','drift_abs','drift_pct','Status']
        if c in drill.columns]
    st.caption(f"{len(drill):,} SKUs")
    st.dataframe(drill[show_cols].rename(columns={
        'v5_first':'1st Tent.','v4':'N-3','v3':'N-2','v2':'N-1','v1_final':'Final',
        'drift_abs':'Δ Units','drift_pct':'Δ %'}).style.format({
        '1st Tent.':'{:,.0f}','N-3':'{:,.0f}','N-2':'{:,.0f}','N-1':'{:,.0f}','Final':'{:,.0f}',
        'Δ Units':'{:+,.0f}','Δ %':'{:+.1f}%'}, na_rep='—'),
    use_container_width=True, hide_index=True, height=480)
    st.download_button("⬇️ Download CSV", drill[show_cols].to_csv(index=False).encode(),
        file_name=f"sku_{selected_month}.csv", mime='text/csv')

with tab5:
    st.markdown("<div class='section-header'>Multi-Month Planning Accuracy</div>", unsafe_allow_html=True)
    if len(all_data) < 2:
        st.info("Upload 2 or more monthly files to see cross-month trends.")
    else:
        months = list(all_data.keys())

        # ── TABLE 1: Month-level summary ──
        st.markdown("**Month Summary**")
        sum_rows = []
        for month, (mdf, _) in all_data.items():
            first = mdf['v5_first'].sum()
            final = mdf['v1_final'].sum()
            sum_rows.append({
                'Month':        month,
                '1st Tentative': first,
                'Final Plan':   final,
                'Drift Units':  final - first,
                'Drift %':      round((final-first)/first*100, 1) if first else 0,
                'Active SKUs':  int(mdf['Is_Active'].sum()),
                'New SKUs':     int((mdf['Status']=='New').sum()),
                'Dropped SKUs': int((mdf['Status']=='Dropped').sum()),
            })
        sum_df = pd.DataFrame(sum_rows)
        sum_styled = sum_df.style.format({'1st Tentative':'{:,.0f}','Final Plan':'{:,.0f}',
                     'Drift Units':'{:+,.0f}','Drift %':'{:+.1f}%'})
        apply_gradient(sum_styled, subset=['Drift %'], col_data=sum_df['Drift %'])
        st.dataframe(sum_styled, use_container_width=True, hide_index=True)

        st.divider()

        # ── CHART: Segment final volumes grouped by month ──
        st.plotly_chart(multi_month_segment_chart(all_data), use_container_width=True)

        st.divider()

        # ── TABLE 2: Segment × Month — Final qty + Drift% per month + cross-period Δ ──
        st.markdown("**Segment Detail — Final Plan & Drift % by Month**")
        all_segs = sorted(set(
            seg for df, _ in all_data.values()
            for seg in df['Segment'].unique()
            if seg not in ('Segment','nan','NaN','','None')
            and (df[df['Segment']==seg]['v1_final'].sum() > 0
                 or df[df['Segment']==seg]['v5_first'].sum() > 0)
        ))
        # Build per-segment finals dict for cross-period calc
        seg_finals_by_month = {}
        seg_rows = []
        for seg in all_segs:
            row = {'Segment': seg}
            seg_finals_by_month[seg] = {}
            for month in months:
                df_m = all_data[month][0]
                first = df_m[df_m['Segment']==seg]['v5_first'].sum()
                final = df_m[df_m['Segment']==seg]['v1_final'].sum()
                dp    = round((final-first)/first*100, 1) if first > 0 else None
                short = month.replace('.xlsx','').replace('.xls','')
                row[f"{short} Final"]  = int(final) if final > 0 else None
                row[f"{short} Drift%"] = dp
                seg_finals_by_month[seg][short] = int(final)
            seg_rows.append(row)

        seg_tbl = pd.DataFrame(seg_rows)
        fmt_seg = {c: '{:,.0f}' for c in seg_tbl.columns if 'Final' in c}
        fmt_seg.update({c: '{:+.1f}%' for c in seg_tbl.columns if 'Drift%' in c})
        drift_cols = [c for c in seg_tbl.columns if 'Drift%' in c]
        styled_seg = seg_tbl.style.format(fmt_seg, na_rep='—')
        if drift_cols:
            for dc in drift_cols:
                apply_gradient(styled_seg, subset=[dc], col_data=seg_tbl[dc])
        st.dataframe(styled_seg, use_container_width=True, hide_index=True)

        st.divider()

        # ── CROSS-PERIOD COMPARISON (Period A vs Period B) ──
        shorts = [m.replace('.xlsx','').replace('.xls','') for m in months]
        if len(months) == 2:
            # Auto-pick A and B
            period_a, period_b = shorts[0], shorts[1]
        else:
            c1, c2 = st.columns(2)
            period_a = c1.selectbox("Base period (A)", shorts, index=0)
            period_b = c2.selectbox("Compare period (B)", shorts, index=len(shorts)-1)

        st.markdown(f"**Period Comparison: {period_a} → {period_b}**")
        st.caption("Shows how final plan volumes shifted for each segment between the two selected periods.")

        cmp_rows = []
        for seg in all_segs:
            fa = seg_finals_by_month[seg].get(period_a, 0)
            fb = seg_finals_by_month[seg].get(period_b, 0)
            delta     = fb - fa
            delta_pct = round(delta / fa * 100, 1) if fa > 0 else None
            cmp_rows.append({
                'Segment':          seg,
                f'{period_a} Final': fa or None,
                f'{period_b} Final': fb or None,
                'Δ Units':           delta if (fa or fb) else None,
                'Δ %':               delta_pct,
            })

        cmp_df = pd.DataFrame(cmp_rows)
        cmp_df = cmp_df[(cmp_df[f'{period_a} Final'].fillna(0) > 0) |
                        (cmp_df[f'{period_b} Final'].fillna(0) > 0)]
        cmp_df = cmp_df.sort_values('Δ %', ascending=False, na_position='last')

        def delta_color(val):
            if pd.isna(val): return 'color: #6b7280'
            if val > 0:  return 'color: #15803d; font-weight:600'
            if val < 0:  return 'color: #dc2626; font-weight:600'
            return 'color: #6b7280'

        cmp_styled = cmp_df.style.format({
                f'{period_a} Final': '{:,.0f}',
                f'{period_b} Final': '{:,.0f}',
                'Δ Units': '{:+,.0f}',
                'Δ %':     '{:+.1f}%',
            }, na_rep='—')
        apply_gradient(cmp_styled, subset=['Δ %'], col_data=cmp_df['Δ %'])
        st.dataframe(cmp_styled, use_container_width=True, hide_index=True)
